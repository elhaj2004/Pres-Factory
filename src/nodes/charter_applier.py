import re
import shutil
import unicodedata
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List

from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor

from src.brand_knowledge import build_brand_quality_targets, get_brand_placeholder_style
from src.reference_styles import resolve_table_profiles
from src.state import PresFactoryState

_ALIGN_DOCX = {
    "left": WD_ALIGN_PARAGRAPH.LEFT,
    "center": WD_ALIGN_PARAGRAPH.CENTER,
    "right": WD_ALIGN_PARAGRAPH.RIGHT,
    "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
}


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _slugify_filename_part(value: str | None) -> str:
    if not value:
        return "document"
    normalized = unicodedata.normalize("NFKD", value)
    ascii_value = normalized.encode("ascii", "ignore").decode("ascii")
    ascii_value = re.sub(r"[^A-Za-z0-9]+", "-", ascii_value).strip("-")
    return ascii_value[:80] or "document"


def _build_output_path(file_path: str, file_type: str, document_title: str | None) -> str:
    output_dir = Path(file_path).parent
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    base_name = _slugify_filename_part(document_title or Path(file_path).stem)
    candidate = output_dir / f"{base_name}-{timestamp}.{file_type}"
    suffix = 2
    while candidate.exists():
        candidate = output_dir / f"{base_name}-{timestamp}-{suffix}.{file_type}"
        suffix += 1
    return str(candidate)


def _clean_text_lines(value: str | None) -> list[str]:
    if not value:
        return []
    parts = re.split(r"[\r\n\v\f]+", value)
    return [" ".join(part.split()).strip() for part in parts if " ".join(part.split()).strip()]


def _extract_source_cover_data(elements: List[Dict[str, Any]]) -> dict[str, Any]:
    cover_elements = [element for element in elements if element.get("slide_idx") == 0]
    title_element = next((element for element in cover_elements if element.get("type") == "title"), None)
    title_lines = _clean_text_lines(title_element.get("content") if title_element else None)
    client_element = next(
        (
            element for element in cover_elements
            if (not title_element or element.get("shape_idx") != title_element.get("shape_idx"))
            and element.get("content")
        ),
        None,
    )
    date_element = next(
        (
            element for element in cover_elements
            if re.fullmatch(r"\d{2}/\d{2}/\d{4}", " ".join((element.get("content") or "").split()))
        ),
        None,
    )
    client_text = " ".join((client_element.get("content") or "").split()).strip() if client_element else ""
    date_text = " ".join((date_element.get("content") or "").split()).strip() if date_element else ""
    return {
        "title_lines": title_lines,
        "client": client_text,
        "date": date_text,
    }


def _to_french_long_date(value: str) -> str:
    match = re.fullmatch(r"(\d{2})/(\d{2})/(\d{4})", value.strip())
    if not match:
        return value
    day, month, year = match.groups()
    month_names = {
        "01": "janvier",
        "02": "fevrier",
        "03": "mars",
        "04": "avril",
        "05": "mai",
        "06": "juin",
        "07": "juillet",
        "08": "aout",
        "09": "septembre",
        "10": "octobre",
        "11": "novembre",
        "12": "decembre",
    }
    return f"{int(day)} {month_names.get(month, month)} {year}"


def _replace_text_preserving_runs(text_frame, paragraphs: list[str]) -> None:
    if not paragraphs:
        return
    text_frame.clear()
    for index, paragraph_text in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = paragraph_text


def _render_from_reference_deck(reference_deck_path: str, elements: List[Dict[str, Any]], output_path: str) -> None:
    shutil.copy2(reference_deck_path, output_path)
    prs = Presentation(output_path)
    cover = _extract_source_cover_data(elements)

    if len(prs.slides) >= 1:
        slide0 = prs.slides[0]
        shapes = list(slide0.shapes)
        title_shape = next((shape for shape in shapes if getattr(shape, "is_placeholder", False) and str(shape.placeholder_format.type) == "CENTER_TITLE (3)"), None)
        body_shape = next((shape for shape in shapes if getattr(shape, "is_placeholder", False) and str(shape.placeholder_format.type) == "BODY (2)"), None)

        title_lines = cover.get("title_lines") or []
        if title_shape and getattr(title_shape, "has_text_frame", False):
            heading = title_lines[0] if title_lines else "Document"
            subtitle = title_lines[1] if len(title_lines) > 1 else ""
            client = cover.get("client") or "Client"
            date = cover.get("date") or ""
            second_line_parts = [part for part in [subtitle, f"[{client}]" if client else "", date] if part]
            title_paragraphs = [heading]
            if second_line_parts:
                title_paragraphs.append(" - ".join(second_line_parts))
            _replace_text_preserving_runs(title_shape.text_frame, title_paragraphs)

        if body_shape and getattr(body_shape, "has_text_frame", False):
            long_date = _to_french_long_date(cover.get("date") or "")
            if long_date:
                _replace_text_preserving_runs(body_shape.text_frame, [long_date])

    prs.save(output_path)


def _preserve_text_content_against_elements(prs: Presentation, elements: List[Dict[str, Any]]) -> None:
    locator_to_text = {
        (element.get("slide_idx"), element.get("shape_idx"), element.get("para_idx")): element.get("content", "")
        for element in elements
        if all(key in element for key in ("slide_idx", "shape_idx", "para_idx"))
    }

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                source_text = locator_to_text.get((slide_idx, shape_idx, para_idx))
                if source_text is None:
                    continue
                current_text = para.text or ""
                if current_text == source_text:
                    continue
                _replace_paragraph_text(para, source_text)


def _replace_paragraph_text(para, new_text: str) -> None:
    runs = list(para.runs)
    if not runs:
        run = para.add_run()
        run.text = new_text
        return
    runs[0].text = new_text
    for run in runs[1:]:
        run.text = ""


def _merge_with_brand_targets(style: Dict[str, Any], placeholder_type: str | None, brand_targets: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(style)
    defaults = get_brand_placeholder_style(placeholder_type, "pptx")
    for key, value in defaults.items():
        if merged.get(key) in (None, ""):
            merged[key] = value

    if merged.get("bold") is None and defaults.get("bold") is not None:
        merged["bold"] = defaults["bold"]
    if merged.get("italic") is None and defaults.get("italic") is not None:
        merged["italic"] = defaults["italic"]

    return merged


def _apply_brand_style_to_shape(shape, shape_style: Dict[str, Any], brand_targets: Dict[str, Any]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    placeholder_type = None
    if getattr(shape, "is_placeholder", False):
        try:
            placeholder_type = str(shape.placeholder_format.type)
        except Exception:
            placeholder_type = "placeholder"

    merged_style = _merge_with_brand_targets(shape_style, placeholder_type, brand_targets)
    for para in shape.text_frame.paragraphs:
        if not para.text.strip():
            continue
        _apply_pptx_para_style(para, merged_style)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _apply_docx_para_style(para, style: Dict[str, Any]) -> None:
    pf = para.paragraph_format
    pf.alignment = _ALIGN_DOCX.get(style.get("alignment", "left"), WD_ALIGN_PARAGRAPH.LEFT)
    if style.get("space_before") is not None:
        pf.space_before = Pt(style["space_before"])
    if style.get("space_after") is not None:
        pf.space_after = Pt(style["space_after"])

    for run in para.runs:
        if style.get("font_name"):
            run.font.name = style["font_name"]
        if style.get("font_size"):
            run.font.size = Pt(float(style["font_size"]))
        if style.get("bold") is not None:
            run.font.bold = style["bold"]
        if style.get("italic") is not None:
            run.font.italic = style["italic"]
        if style.get("color"):
            r, g, b = _hex_to_rgb(style["color"])
            run.font.color.rgb = RGBColor(r, g, b)


def _apply_docx(file_path: str, elements: List[Dict], style_map: List[Dict], output_path: str) -> None:
    doc = DocxDocument(file_path)

    # Index styles by element id
    style_by_id: Dict[str, Dict] = {s["id"]: s for s in style_map if "id" in s}

    # Match paragraphs in order (parser et applier utilisent le même ordre)
    elem_iter = iter(elements)
    current_elem = next(elem_iter, None)

    for para in doc.paragraphs:
        if not para.text.strip() or current_elem is None:
            continue
        style = style_by_id.get(current_elem["id"])
        if style:
            _apply_docx_para_style(para, style)
        current_elem = next(elem_iter, None)

    header_style, body_style = resolve_table_profiles(style_map, "docx", {})
    for table in doc.tables:
        for row_idx, row in enumerate(table.rows):
            row_style = header_style if row_idx == 0 and header_style else body_style
            if not row_style:
                continue
            for cell in row.cells:
                for cell_para in cell.paragraphs:
                    if row_style.get("alignment"):
                        cell_para.paragraph_format.alignment = _ALIGN_DOCX.get(
                            row_style.get("alignment", "left"),
                            WD_ALIGN_PARAGRAPH.LEFT,
                        )
                    for run in cell_para.runs:
                        if row_style.get("font_name"):
                            run.font.name = row_style["font_name"]
                        if row_style.get("font_size"):
                            run.font.size = Pt(float(row_style["font_size"]))
                        if row_style.get("bold") is not None:
                            run.font.bold = row_style.get("bold")
                        if row_style.get("italic") is not None:
                            run.font.italic = row_style.get("italic")
                        if row_style.get("color"):
                            r, g, b = _hex_to_rgb(row_style["color"])
                            run.font.color.rgb = RGBColor(r, g, b)

    doc.save(output_path)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _apply_pptx_para_style(para, style: Dict[str, Any]) -> None:
    if not para.runs and para.text:
        para.add_run()

    if style.get("alignment"):
        align_map = {
            "left": 1,
            "center": 2,
            "right": 3,
            "justify": 4,
        }
        mapped = align_map.get(style.get("alignment"))
        if mapped is not None:
            para.alignment = mapped

    for run in para.runs:
        if style.get("font_name"):
            run.font.name = style["font_name"]
        if style.get("font_size"):
            run.font.size = PptPt(float(style["font_size"]))
        if style.get("bold") is not None:
            run.font.bold = style["bold"]
        if style.get("italic") is not None:
            run.font.italic = style["italic"]
        if style.get("color"):
            r, g, b = _hex_to_rgb(style["color"])
            run.font.color.rgb = PptRGBColor(r, g, b)


def _apply_pptx(file_path: str, elements: List[Dict], style_map: List[Dict], output_path: str, brand_targets: Dict[str, Any]) -> None:
    prs = Presentation(file_path)
    style_by_id: Dict[str, Dict] = {s["id"]: s for s in style_map if "id" in s}

    shape_style_by_locator: Dict[tuple[int, int], Dict[str, Any]] = {}
    for style in style_map:
        slide_idx = style.get("slide_idx")
        shape_idx = style.get("shape_idx")
        if slide_idx is None or shape_idx is None:
            continue
        key = (slide_idx, shape_idx)
        if key in shape_style_by_locator:
            continue
        shape_style_by_locator[key] = style

    locator_to_id = {
        (element.get("slide_idx"), element.get("shape_idx"), element.get("para_idx")): element["id"]
        for element in elements
        if all(key in element for key in ("slide_idx", "shape_idx", "para_idx"))
    }

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue

            shape_style = shape_style_by_locator.get((slide_idx, shape_idx))
            if shape_style and getattr(shape.text_frame, "paragraphs", None):
                _apply_brand_style_to_shape(shape, shape_style, brand_targets)

            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                element_id = locator_to_id.get((slide_idx, shape_idx, para_idx))
                if not element_id:
                    if shape_style and para.text.strip():
                        placeholder_type = None
                        if getattr(shape, "is_placeholder", False):
                            try:
                                placeholder_type = str(shape.placeholder_format.type)
                            except Exception:
                                placeholder_type = "placeholder"
                        _apply_pptx_para_style(para, _merge_with_brand_targets(shape_style, placeholder_type, brand_targets))
                    continue
                style = style_by_id.get(element_id)
                if style:
                    placeholder_type = None
                    if getattr(shape, "is_placeholder", False):
                        try:
                            placeholder_type = str(shape.placeholder_format.type)
                        except Exception:
                            placeholder_type = "placeholder"
                    _apply_pptx_para_style(para, _merge_with_brand_targets(style, placeholder_type, brand_targets))

    _preserve_text_content_against_elements(prs, elements)

    prs.save(output_path)


# ── NODE ──────────────────────────────────────────────────────────────────────

def apply_charter(state: PresFactoryState) -> dict:
    file_path = state["file_path"]
    file_type = state["file_type"]
    elements = state["anonymized_elements"]
    style_map = state["style_map"]
    primary_reference_deck = state.get("primary_reference_deck") or {}
    brand_targets = build_brand_quality_targets(file_type)
    iteration = state.get("iteration_count", 0)
    document_title = state.get("document_title")

    output_path = _build_output_path(file_path, file_type, document_title)

    try:
        if file_type == "pptx" and primary_reference_deck.get("strategy") == "copy_reference_deck":
            _render_from_reference_deck(primary_reference_deck["path"], elements, output_path)
        elif file_type == "docx":
            _apply_docx(file_path, elements, style_map, output_path)
        else:
            _apply_pptx(file_path, elements, style_map, output_path, brand_targets)

        return {
            "output_path": output_path,
            "iteration_count": iteration + 1,
        }
    except Exception as e:
        return {"error": f"Erreur application charte: {e}"}
