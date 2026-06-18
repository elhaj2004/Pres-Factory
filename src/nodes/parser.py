import json
import re
from pathlib import Path
from typing import List, Dict, Any

from docx import Document as DocxDocument
from pptx import Presentation

from src.state import PresFactoryState


def _shape_geometry_bucket(shape) -> str:
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except Exception:
        return ""

    horizontal = "left" if left < 2000000 else ("right" if left > 5000000 else "center")
    vertical = "top" if top < 1200000 else ("bottom" if top > 3000000 else "middle")
    size = "large" if width > 5000000 or height > 1200000 else "small"
    return f"{vertical}_{horizontal}_{size}"


def _docx_element_type(style_name: str) -> tuple[str, int]:
    name = style_name.lower()
    if "heading 1" in name:
        return "heading_1", 1
    if "heading 2" in name:
        return "heading_2", 2
    if "heading 3" in name:
        return "heading_3", 3
    if any(k in name for k in ("list", "bullet", "liste")):
        return "bullet", 1
    if "caption" in name or "légende" in name:
        return "caption", 0
    return "body", 0


def _parse_docx(file_path: str) -> List[Dict[str, Any]]:
    doc = DocxDocument(file_path)
    elements = []
    idx = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        elem_type, level = _docx_element_type(para.style.name)
        elements.append({
            "id": f"para_{idx}",
            "type": elem_type,
            "content": text,
            "original_style": para.style.name,
            "level": level,
            "source": "paragraph",
        })
        idx += 1

    for t_idx, table in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        elements.append({
            "id": f"table_{t_idx}",
            "type": "table",
            "content": json.dumps(rows, ensure_ascii=False),
            "original_style": "table",
            "level": 0,
            "source": "table",
        })

    return elements


def _pptx_element_type(shape_idx: int, para_level: int, is_first_shape: bool) -> str:
    if is_first_shape and shape_idx == 0:
        return "title"
    if para_level == 0:
        return "heading" if shape_idx == 0 else "body"
    if para_level == 1:
        return "bullet_level_1"
    return "bullet_level_2"


def _parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    prs = Presentation(file_path)
    elements = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            placeholder_type = None
            if getattr(shape, "is_placeholder", False):
                try:
                    placeholder_type = str(shape.placeholder_format.type)
                except Exception:
                    placeholder_type = "placeholder"
            geometry_bucket = _shape_geometry_bucket(shape)
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                elem_type = _pptx_element_type(
                    shape_idx, para.level, slide_idx == 0
                )
                elements.append({
                    "id": f"s{slide_idx}_sh{shape_idx}_p{para_idx}",
                    "type": elem_type,
                    "content": text,
                    "original_style": None,
                    "level": para.level,
                    "source": "slide",
                    "slide_idx": slide_idx,
                    "shape_idx": shape_idx,
                    "para_idx": para_idx,
                    "placeholder_type": placeholder_type,
                    "shape_name": getattr(shape, "name", None),
                    "geometry_bucket": geometry_bucket,
                })

    return elements


def _derive_document_title(file_path: str, elements: List[Dict[str, Any]]) -> str | None:
    def _clean_title_candidate(raw_text: str) -> str:
        normalized = re.split(r"[\r\n\v\f]+", raw_text or "")
        first_line = next((part.strip() for part in normalized if part.strip()), "")
        return " ".join(first_line.split()).strip()

    preferred_types = ["title", "heading_1", "heading", "heading_2", "body"]

    for elem_type in preferred_types:
        for element in elements:
            if element.get("type") != elem_type:
                continue
            text = _clean_title_candidate(element.get("content") or "")
            if len(text) >= 4:
                return text

    fallback = Path(file_path).stem.strip()
    return fallback or None


def parse_document(state: PresFactoryState) -> dict:
    try:
        if state["file_type"] == "docx":
            elements = _parse_docx(state["file_path"])
        else:
            elements = _parse_pptx(state["file_path"])
        return {
            "raw_elements": elements,
            "document_title": _derive_document_title(state["file_path"], elements),
        }
    except Exception as e:
        return {"error": f"Erreur de parsing: {e}"}
