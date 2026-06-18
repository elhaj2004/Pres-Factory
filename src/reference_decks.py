import re
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from src.brand_knowledge import get_official_template_decks


def _tokenize(text: str) -> list[str]:
    return [token for token in re.findall(r"[a-z0-9]+", text.lower()) if len(token) > 2]


def _normalize_title_tokens(value: str | None) -> set[str]:
    if not value:
        return set()
    return set(_tokenize(value))


def _collect_pptx_summary(path: Path) -> dict[str, Any] | None:
    if not path.exists() or path.suffix.lower() != ".pptx":
        return None

    try:
        prs = Presentation(str(path))
    except Exception:
        return None

    slide_texts: list[str] = []
    layout_names: list[str] = []
    for slide in prs.slides:
        layout_names.append(getattr(slide.slide_layout, "name", "") or "")
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                parts.append(" ".join(text.split()))
        slide_texts.append(" ".join(parts))

    full_text = "\n".join(slide_texts)
    tokens = Counter(_tokenize(full_text))
    return {
        "path": str(path),
        "slide_count": len(slide_texts),
        "layout_variety": len({name for name in layout_names if name}),
        "layout_names": layout_names,
        "slide_texts": slide_texts,
        "tokens": tokens,
        "token_total": sum(tokens.values()),
        "title_tokens": _normalize_title_tokens(path.stem),
    }


def _collect_target_summary(file_path: str, document_title: str | None, raw_elements: list[dict[str, Any]]) -> dict[str, Any] | None:
    path = Path(file_path)
    if path.suffix.lower() != ".pptx":
        return None

    slide_texts: dict[int, list[str]] = {}
    for element in raw_elements or []:
        if not isinstance(element, dict) or "slide_idx" not in element:
            continue
        text = " ".join((element.get("content") or "").split()).strip()
        if not text:
            continue
        slide_texts.setdefault(int(element["slide_idx"]), []).append(text)

    ordered_slide_texts = [" ".join(slide_texts[idx]) for idx in sorted(slide_texts)]
    tokens = Counter(_tokenize("\n".join(ordered_slide_texts)))
    return {
        "path": str(path),
        "slide_count": len(ordered_slide_texts),
        "layout_variety": 0,
        "slide_texts": ordered_slide_texts,
        "tokens": tokens,
        "token_total": sum(tokens.values()),
        "title_tokens": _normalize_title_tokens(document_title or path.stem),
    }


def find_primary_reference_deck(
    original_file_path: str | None,
    file_type: str | None,
    document_title: str | None,
    raw_elements: list[dict[str, Any]],
) -> dict[str, Any] | None:
    if file_type != "pptx" or not original_file_path:
        return None

    original_path = Path(original_file_path)
    if not original_path.exists() or original_path.suffix.lower() != ".pptx":
        return None

    target = _collect_target_summary(str(original_path), document_title, raw_elements)
    if not target:
        return None

    parent = original_path.parent
    official_template_paths = [Path(path) for path in get_official_template_decks(file_type or "")]
    candidates: list[dict[str, Any]] = []
    candidate_pool = list(parent.glob("*.pptx"))
    for official_path in official_template_paths:
        if official_path.exists():
            candidate_pool.append(official_path)

    seen_paths: set[str] = set()
    for candidate_path in candidate_pool:
        resolved = str(candidate_path.resolve())
        if resolved in seen_paths:
            continue
        seen_paths.add(resolved)
        if candidate_path.resolve() == original_path.resolve():
            continue
        candidate = _collect_pptx_summary(candidate_path)
        if not candidate:
            continue

        title_overlap = len(target["title_tokens"] & candidate["title_tokens"])
        overlap_count = sum((target["tokens"] & candidate["tokens"]).values())
        overlap_ratio = overlap_count / max(target["token_total"], 1)
        slide_delta = candidate["slide_count"] - target["slide_count"]
        layout_bonus = 20 if candidate["layout_variety"] >= 3 else 0
        structured_bonus = 60 if slide_delta >= 3 else 0
        same_size_penalty = 35 if slide_delta == 0 else 0
        title_bonus = title_overlap * 30
        is_official = candidate_path in official_template_paths
        official_bonus = 140 if is_official else 0
        score = overlap_ratio * 100 + structured_bonus + layout_bonus + title_bonus + official_bonus - same_size_penalty

        candidates.append(
            {
                **candidate,
                "title_overlap": title_overlap,
                "overlap_count": overlap_count,
                "overlap_ratio": round(overlap_ratio, 4),
                "score": round(score, 2),
                "slide_delta": slide_delta,
                "is_official_template": is_official,
            }
        )

    if not candidates:
        return None

    candidates.sort(
        key=lambda item: (
            item["score"],
            item["slide_delta"],
            item["title_overlap"],
            item["overlap_count"],
        ),
        reverse=True,
    )
    best = candidates[0]
    if best["title_overlap"] == 0 or best["score"] < 80:
        return None

    strategy = "copy_reference_deck" if best["slide_delta"] >= 3 else "style_only"
    return {
        "path": best["path"],
        "strategy": strategy,
        "score": best["score"],
        "title_overlap": best["title_overlap"],
        "overlap_ratio": best["overlap_ratio"],
        "slide_delta": best["slide_delta"],
        "is_official_template": best.get("is_official_template", False),
    }
