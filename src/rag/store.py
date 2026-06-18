import hashlib
import json
import os
import shutil
from dataclasses import dataclass
from pathlib import Path, PurePosixPath
from typing import Any, Optional
from urllib.parse import parse_qs, quote, unquote, urlparse

from docx import Document as WordDocument
from dotenv import load_dotenv
from langchain_community.vectorstores import Chroma
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

from src.llm.client import get_embeddings

load_dotenv()

PROJECT_ROOT = Path(__file__).resolve().parents[2]
DATA_DIR = PROJECT_ROOT / "data"
DEFAULT_LOCAL_BRANDSTORE_DIR = DATA_DIR / "brandstore"
CHROMA_ROOT_DIR = DATA_DIR / "chroma_db"
SHAREPOINT_CACHE_ROOT_DIR = DATA_DIR / "sharepoint_cache"
DEFAULT_SHAREPOINT_FOLDER_URL = (
    "https://orangecyberdefense.sharepoint.com/sites/directionExpertise/"
    "Documents%20partages/Forms/AllItems.aspx"
)
GRAPH_API_BASE_URL = "https://graph.microsoft.com/v1.0"
SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".txt", ".md"}
TEXT_SEPARATORS = ["\n\n", "\n", ". ", " "]

_store: Optional[Chroma] = None
_store_key: Optional[str] = None


class SharePointSyncError(RuntimeError):
    pass


class NoOpEmbeddings(Embeddings):
    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        return [[0.0] for _ in texts]

    def embed_query(self, text: str) -> list[float]:
        return [0.0]


@dataclass(frozen=True)
class SharePointTarget:
    host: str
    site_path: str
    drive_name: str
    folder_path: str
    source_url: str


@dataclass(frozen=True)
class RagSourceConfig:
    source_type: str
    source_id: str
    source_label: str
    documents_dir: Path
    persist_dir: Path
    collection_name: str
    signature: str
    manifest_path: Optional[Path] = None

    @property
    def index_state_path(self) -> Path:
        return self.persist_dir / "index_state.json"

    @property
    def store_cache_key(self) -> str:
        return f"{self.source_id}:{self.signature}"


def get_rag_store() -> Chroma:
    global _store, _store_key

    config = _prepare_source(force_refresh=False, verbose=False)
    index_state = _load_json(config.index_state_path)
    if _store is not None and _store_key == config.store_cache_key:
        return _store

    if _should_reindex(config):
        index_brandstore(verbose=False)
        if _store is not None and _store_key == config.store_cache_key:
            return _store

    config.persist_dir.mkdir(parents=True, exist_ok=True)
    _store = Chroma(
        persist_directory=str(config.persist_dir),
        embedding_function=_get_embeddings_for_store(index_state),
        collection_name=config.collection_name,
    )
    _store_key = config.store_cache_key
    return _store


def index_brandstore(verbose: bool = True, force_refresh: bool = False) -> int:
    config = _prepare_source(force_refresh=force_refresh, verbose=verbose)
    docs = _load_documents_from_directory(config.documents_dir, config.source_type)

    global _store, _store_key
    _store = None
    _store_key = None

    if config.persist_dir.exists():
        shutil.rmtree(config.persist_dir)
    config.persist_dir.mkdir(parents=True, exist_ok=True)

    if not docs:
        if verbose:
            print(
                f"[RAG] Aucun document supporte a indexer pour la source {config.source_label}"
            )
        store = _build_empty_store(config)
        _write_index_state(config, document_count=0, chunk_count=0)
        _store = store
        _store_key = config.store_cache_key
        return 0

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        separators=TEXT_SEPARATORS,
    )
    chunks = splitter.split_documents(docs)

    store = Chroma.from_documents(
        documents=chunks,
        embedding=get_embeddings(),
        persist_directory=str(config.persist_dir),
        collection_name=config.collection_name,
    )
    _persist_store(store)
    _write_index_state(config, document_count=len(docs), chunk_count=len(chunks))

    if verbose:
        print(
            f"[RAG] Source={config.source_label} | documents={len(docs)} | chunks={len(chunks)}"
        )

    _store = store
    _store_key = config.store_cache_key
    return len(chunks)


def similarity_search(query: str, k: int = 3) -> list:
    try:
        return get_rag_store().similarity_search(query, k=k)
    except Exception:
        return []


def _prepare_source(force_refresh: bool, verbose: bool) -> RagSourceConfig:
    source = os.getenv("RAG_SOURCE", "local").strip().lower() or "local"
    if source == "sharepoint":
        return _prepare_sharepoint_source(force_refresh=force_refresh, verbose=verbose)
    if source != "local":
        raise ValueError(f"RAG_SOURCE invalide: {source}. Valeurs attendues: local, sharepoint")
    return _prepare_local_source()


def _prepare_local_source() -> RagSourceConfig:
    documents_dir = _resolve_path(
        os.getenv("BRANDSTORE_DIR"),
        DEFAULT_LOCAL_BRANDSTORE_DIR,
    )
    source_hash = hashlib.sha1(str(documents_dir.resolve()).encode("utf-8")).hexdigest()[:12]
    return RagSourceConfig(
        source_type="local",
        source_id=f"local_{source_hash}",
        source_label=f"local:{documents_dir}",
        documents_dir=documents_dir,
        persist_dir=CHROMA_ROOT_DIR / f"local_{source_hash}",
        collection_name=f"brandstore_{source_hash}",
        signature=_compute_directory_signature(documents_dir),
    )


def _prepare_sharepoint_source(force_refresh: bool, verbose: bool) -> RagSourceConfig:
    target = _resolve_sharepoint_target()
    source_key = json.dumps(
        {
            "host": target.host,
            "site_path": target.site_path,
            "drive_name": target.drive_name,
            "folder_path": target.folder_path,
        },
        sort_keys=True,
    )
    source_hash = hashlib.sha1(source_key.encode("utf-8")).hexdigest()[:12]
    cache_root = SHAREPOINT_CACHE_ROOT_DIR / f"sharepoint_{source_hash}"
    documents_dir = _sync_sharepoint_cache(
        target=target,
        cache_root=cache_root,
        force_refresh=force_refresh,
        verbose=verbose,
    )
    manifest_path = cache_root / "manifest.json"
    return RagSourceConfig(
        source_type="sharepoint",
        source_id=f"sharepoint_{source_hash}",
        source_label=(
            f"sharepoint:{target.host}{target.site_path}/{target.drive_name}/{target.folder_path}"
            .rstrip("/")
        ),
        documents_dir=documents_dir,
        persist_dir=CHROMA_ROOT_DIR / f"sharepoint_{source_hash}",
        collection_name=f"brandstore_{source_hash}",
        signature=_compute_sharepoint_signature(manifest_path, documents_dir),
        manifest_path=manifest_path,
    )


def _resolve_sharepoint_target() -> SharePointTarget:
    folder_url = (
        os.getenv("SHAREPOINT_FOLDER_URL")
        or os.getenv("SHAREPOINT_SITE_URL")
        or DEFAULT_SHAREPOINT_FOLDER_URL
    )
    target = _parse_sharepoint_folder_url(folder_url)

    site_url = os.getenv("SHAREPOINT_SITE_URL")
    if site_url:
        host, site_path = _parse_sharepoint_site_url(site_url)
        target = SharePointTarget(
            host=host,
            site_path=site_path,
            drive_name=target.drive_name,
            folder_path=target.folder_path,
            source_url=target.source_url,
        )

    drive_name = os.getenv("SHAREPOINT_DOCUMENT_LIBRARY") or target.drive_name
    folder_path = os.getenv("SHAREPOINT_FOLDER_PATH")
    if folder_path is None:
        folder_path = target.folder_path

    folder_path = _normalize_sharepoint_folder_path(folder_path, drive_name)
    return SharePointTarget(
        host=target.host,
        site_path=target.site_path,
        drive_name=drive_name,
        folder_path=folder_path,
        source_url=target.source_url,
    )


def _parse_sharepoint_site_url(site_url: str) -> tuple[str, str]:
    parsed = urlparse(site_url)
    if not parsed.scheme or not parsed.netloc:
        raise ValueError(f"URL SharePoint invalide: {site_url}")

    segments = [segment for segment in unquote(parsed.path).split("/") if segment]
    if len(segments) < 2 or segments[0] not in {"sites", "teams"}:
        raise ValueError(
            "Impossible de determiner le site SharePoint. URL attendue: /sites/<nom> ou /teams/<nom>."
        )
    return parsed.netloc, f"/{segments[0]}/{segments[1]}"


def _parse_sharepoint_folder_url(folder_url: str) -> SharePointTarget:
    host, site_path = _parse_sharepoint_site_url(folder_url)
    parsed = urlparse(folder_url)
    path_segments = [segment for segment in unquote(parsed.path).split("/") if segment]
    query = parse_qs(parsed.query)

    remaining_segments = path_segments[2:]
    if "id" in query and query["id"]:
        id_segments = [segment for segment in unquote(query["id"][0]).split("/") if segment]
        if len(id_segments) >= 3 and id_segments[0:2] == path_segments[0:2]:
            remaining_segments = id_segments[2:]

    cleaned_segments: list[str] = []
    for segment in remaining_segments:
        if segment.lower() == "forms":
            break
        cleaned_segments.append(segment)

    drive_name = cleaned_segments[0] if cleaned_segments else "Documents partages"
    folder_path = "/".join(cleaned_segments[1:]) if len(cleaned_segments) > 1 else ""
    return SharePointTarget(
        host=host,
        site_path=site_path,
        drive_name=drive_name,
        folder_path=folder_path,
        source_url=folder_url,
    )


def _normalize_sharepoint_folder_path(folder_path: str, drive_name: str) -> str:
    normalized = folder_path.strip().strip("/")
    if not normalized:
        return ""
    if normalized == drive_name:
        return ""
    drive_prefix = f"{drive_name}/"
    if normalized.startswith(drive_prefix):
        return normalized[len(drive_prefix) :]
    return normalized


def _sync_sharepoint_cache(
    target: SharePointTarget,
    cache_root: Path,
    force_refresh: bool,
    verbose: bool,
) -> Path:
    files_dir = cache_root / "files"
    manifest_path = cache_root / "manifest.json"
    files_dir.mkdir(parents=True, exist_ok=True)
    previous_manifest = _load_json(manifest_path)

    try:
        session = _create_graph_session()
        site = _graph_get_json(
            session,
            f"{GRAPH_API_BASE_URL}/sites/{target.host}:{target.site_path}?$select=id,webUrl",
        )
        drive = _find_drive(session, site["id"], target.drive_name)
        folder_item = _get_drive_folder(session, drive["id"], target.folder_path)
        remote_files = _list_supported_drive_files(
            session,
            drive_id=drive["id"],
            folder_item_id=folder_item["id"],
        )

        remote_manifest: dict[str, dict[str, Any]] = {}
        for remote_file in remote_files:
            file_state = {
                "id": remote_file["id"],
                "name": remote_file["name"],
                "relative_path": remote_file["relative_path"],
                "last_modified": remote_file.get("lastModifiedDateTime", ""),
                "size": remote_file.get("size", 0),
                "etag": remote_file.get("eTag", ""),
            }
            remote_manifest[remote_file["id"]] = file_state

            local_path = _cache_path(files_dir, remote_file["relative_path"])
            local_path.parent.mkdir(parents=True, exist_ok=True)

            previous_state = previous_manifest.get("files", {}).get(remote_file["id"], {})
            previous_relative_path = previous_state.get("relative_path")
            if previous_relative_path and previous_relative_path != remote_file["relative_path"]:
                previous_path = _cache_path(files_dir, previous_relative_path)
                if previous_path.exists():
                    previous_path.unlink()

            unchanged = (
                not force_refresh
                and previous_state.get("last_modified") == file_state["last_modified"]
                and previous_state.get("size") == file_state["size"]
                and previous_state.get("etag") == file_state["etag"]
                and local_path.exists()
            )
            if unchanged:
                continue

            _download_drive_item(session, drive["id"], remote_file["id"], local_path)

        previous_files = previous_manifest.get("files", {}) if isinstance(previous_manifest, dict) else {}
        current_ids = set(remote_manifest)
        for previous_id, previous_state in previous_files.items():
            if previous_id in current_ids:
                continue
            previous_relative_path = previous_state.get("relative_path")
            if not previous_relative_path:
                continue
            stale_path = _cache_path(files_dir, previous_relative_path)
            if stale_path.exists():
                stale_path.unlink()

        _cleanup_empty_directories(files_dir)
        _write_json(
            manifest_path,
            {
                "source_url": target.source_url,
                "host": target.host,
                "site_path": target.site_path,
                "drive_name": target.drive_name,
                "folder_path": target.folder_path,
                "site_id": site["id"],
                "drive_id": drive["id"],
                "folder_item_id": folder_item["id"],
                "files": remote_manifest,
            },
        )

        if verbose:
            print(
                f"[RAG] SharePoint synchronise: {len(remote_manifest)} fichiers supportes vers {files_dir}"
            )
        return files_dir
    except SharePointSyncError as exc:
        if _has_usable_cache(files_dir):
            if verbose:
                print(f"[RAG] Sync SharePoint indisponible ({exc}). Utilisation du cache local.")
            return files_dir
        raise


def _create_graph_session():
    tenant_id = os.getenv("SHAREPOINT_TENANT_ID", "").strip()
    client_id = os.getenv("SHAREPOINT_CLIENT_ID", "").strip()
    client_secret = os.getenv("SHAREPOINT_CLIENT_SECRET", "").strip()
    if not tenant_id or not client_id or not client_secret:
        raise SharePointSyncError(
            "Credentials Microsoft Graph manquants (SHAREPOINT_TENANT_ID, SHAREPOINT_CLIENT_ID, SHAREPOINT_CLIENT_SECRET)."
        )

    try:
        import msal
        import requests
    except ImportError as exc:
        raise SharePointSyncError(f"Dependance SharePoint manquante: {exc}") from exc

    app = msal.ConfidentialClientApplication(
        client_id=client_id,
        authority=f"https://login.microsoftonline.com/{tenant_id}",
        client_credential=client_secret,
    )
    try:
        token_result = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    except Exception as exc:
        raise SharePointSyncError(f"Echec d'authentification Graph: {exc}") from exc
    access_token = token_result.get("access_token")
    if not access_token:
        description = token_result.get("error_description") or token_result.get("error") or "echec inconnu"
        raise SharePointSyncError(f"Echec d'authentification Graph: {description}")

    session = requests.Session()
    session.headers.update({"Authorization": f"Bearer {access_token}"})
    return session


def _find_drive(session, site_id: str, drive_name: str) -> dict[str, Any]:
    drives = _graph_get_json(
        session,
        f"{GRAPH_API_BASE_URL}/sites/{site_id}/drives?$select=id,name,driveType,webUrl",
    ).get("value", [])
    for drive in drives:
        if drive.get("name", "").casefold() == drive_name.casefold():
            return drive
    raise SharePointSyncError(f"Bibliotheque SharePoint introuvable: {drive_name}")


def _get_drive_folder(session, drive_id: str, folder_path: str) -> dict[str, Any]:
    if not folder_path:
        return _graph_get_json(
            session,
            f"{GRAPH_API_BASE_URL}/drives/{drive_id}/root?$select=id,name,webUrl",
        )

    encoded_path = quote(folder_path.strip("/"), safe="/")
    return _graph_get_json(
        session,
        f"{GRAPH_API_BASE_URL}/drives/{drive_id}/root:/{encoded_path}:?$select=id,name,webUrl",
    )


def _list_supported_drive_files(session, drive_id: str, folder_item_id: str) -> list[dict[str, Any]]:
    files: list[dict[str, Any]] = []
    stack: list[tuple[str, PurePosixPath]] = [(folder_item_id, PurePosixPath("."))]

    while stack:
        item_id, base_path = stack.pop()
        for child in _iterate_drive_children(session, drive_id, item_id):
            child_name = child.get("name", "")
            child_path = base_path / child_name
            if child.get("folder"):
                stack.append((child["id"], child_path))
                continue
            if not child.get("file"):
                continue
            if Path(child_name).suffix.lower() not in SUPPORTED_EXTENSIONS:
                continue
            files.append(
                {
                    **child,
                    "relative_path": (
                        child_name
                        if base_path == PurePosixPath(".")
                        else child_path.as_posix()
                    ),
                }
            )

    files.sort(key=lambda item: item["relative_path"])
    return files


def _iterate_drive_children(session, drive_id: str, item_id: str):
    next_url = (
        f"{GRAPH_API_BASE_URL}/drives/{drive_id}/items/{item_id}/children"
        "?$select=id,name,eTag,size,lastModifiedDateTime,file,folder"
    )
    while next_url:
        payload = _graph_get_json(session, next_url)
        for item in payload.get("value", []):
            yield item
        next_url = payload.get("@odata.nextLink")


def _download_drive_item(session, drive_id: str, item_id: str, destination: Path) -> None:
    try:
        response = session.get(
            f"{GRAPH_API_BASE_URL}/drives/{drive_id}/items/{item_id}/content",
            timeout=120,
            stream=True,
        )
    except Exception as exc:
        raise SharePointSyncError(f"Telechargement impossible pour l'item {item_id}: {exc}") from exc
    if response.status_code >= 400:
        detail = _extract_graph_error(response)
        raise SharePointSyncError(f"Telechargement impossible pour l'item {item_id}: {detail}")

    with destination.open("wb") as handle:
        for chunk in response.iter_content(chunk_size=1024 * 1024):
            if chunk:
                handle.write(chunk)


def _graph_get_json(session, url: str) -> dict[str, Any]:
    try:
        response = session.get(url, timeout=60)
    except Exception as exc:
        raise SharePointSyncError(f"Appel Graph impossible: {exc}") from exc
    if response.status_code >= 400:
        detail = _extract_graph_error(response)
        raise SharePointSyncError(detail)
    return response.json()


def _extract_graph_error(response) -> str:
    try:
        payload = response.json()
    except ValueError:
        return f"HTTP {response.status_code}: {response.text.strip()}"

    error = payload.get("error") if isinstance(payload, dict) else None
    if isinstance(error, dict):
        message = error.get("message") or error.get("code") or str(error)
        return f"HTTP {response.status_code}: {message}"
    return f"HTTP {response.status_code}: {payload}"


def _load_documents_from_directory(documents_dir: Path, source_type: str) -> list[Document]:
    if not documents_dir.exists():
        return []

    docs: list[Document] = []
    for file_path in sorted(documents_dir.rglob("*")):
        if not file_path.is_file() or file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        content = _extract_text(file_path)
        if not content.strip():
            continue

        try:
            relative_path = file_path.relative_to(documents_dir).as_posix()
        except ValueError:
            relative_path = file_path.name

        docs.append(
            Document(
                page_content=content,
                metadata={
                    "source": str(file_path),
                    "relative_path": relative_path,
                    "filename": file_path.name,
                    "extension": file_path.suffix.lower(),
                    "rag_source": source_type,
                },
            )
        )
    return docs


def _extract_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        return _extract_docx_text(file_path)
    if suffix == ".pptx":
        return _extract_pptx_text(file_path)
    if suffix in {".txt", ".md"}:
        return _extract_text_file(file_path)
    return ""


def _extract_docx_text(file_path: Path) -> str:
    try:
        document = WordDocument(str(file_path))
    except Exception:
        return ""

    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(text)

    for table in document.tables:
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_values:
                blocks.append(" | ".join(row_values))
    return "\n".join(blocks)


def _extract_pptx_text(file_path: Path) -> str:
    try:
        presentation = Presentation(str(file_path))
    except Exception:
        return ""

    slides_text: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        blocks = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                blocks.append(text.strip())
                continue
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_values:
                        blocks.append(" | ".join(row_values))
        if len(blocks) > 1:
            slides_text.append("\n".join(blocks))
    return "\n\n".join(slides_text)


def _extract_text_file(file_path: Path) -> str:
    data = file_path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _build_empty_store(config: RagSourceConfig) -> Chroma:
    store = Chroma(
        persist_directory=str(config.persist_dir),
        embedding_function=NoOpEmbeddings(),
        collection_name=config.collection_name,
    )
    _persist_store(store)
    return store


def _get_embeddings_for_store(index_state: dict[str, Any]) -> Embeddings:
    if index_state.get("document_count", 0) == 0:
        return NoOpEmbeddings()
    return get_embeddings()


def _persist_store(store: Chroma) -> None:
    persist = getattr(store, "persist", None)
    if callable(persist):
        persist()


def _should_reindex(config: RagSourceConfig) -> bool:
    if not config.persist_dir.exists():
        return True
    index_state = _load_json(config.index_state_path)
    return index_state.get("source_id") != config.source_id or index_state.get("source_signature") != config.signature


def _write_index_state(config: RagSourceConfig, document_count: int, chunk_count: int) -> None:
    _write_json(
        config.index_state_path,
        {
            "source_id": config.source_id,
            "source_type": config.source_type,
            "source_label": config.source_label,
            "source_signature": config.signature,
            "document_count": document_count,
            "chunk_count": chunk_count,
        },
    )


def _compute_directory_signature(directory: Path) -> str:
    digest = hashlib.sha256()
    if not directory.exists():
        digest.update(b"missing")
        return digest.hexdigest()

    for file_path in sorted(directory.rglob("*")):
        if not file_path.is_file() or file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        relative_path = file_path.relative_to(directory).as_posix()
        stats = file_path.stat()
        digest.update(relative_path.encode("utf-8"))
        digest.update(str(stats.st_size).encode("utf-8"))
        digest.update(str(stats.st_mtime_ns).encode("utf-8"))
    return digest.hexdigest()


def _compute_sharepoint_signature(manifest_path: Path, documents_dir: Path) -> str:
    manifest = _load_json(manifest_path)
    files = manifest.get("files") if isinstance(manifest, dict) else None
    if isinstance(files, dict):
        digest = hashlib.sha256()
        for item_id in sorted(files):
            file_state = files[item_id]
            digest.update(item_id.encode("utf-8"))
            digest.update(file_state.get("relative_path", "").encode("utf-8"))
            digest.update(str(file_state.get("size", 0)).encode("utf-8"))
            digest.update(file_state.get("last_modified", "").encode("utf-8"))
            digest.update(file_state.get("etag", "").encode("utf-8"))
        return digest.hexdigest()
    return _compute_directory_signature(documents_dir)


def _resolve_path(raw_path: Optional[str], default_path: Path) -> Path:
    if not raw_path:
        return default_path
    path = Path(raw_path)
    if path.is_absolute():
        return path
    return PROJECT_ROOT / path


def _cache_path(files_dir: Path, relative_path: str) -> Path:
    normalized_path = PurePosixPath(relative_path)
    if any(part == ".." for part in normalized_path.parts):
        raise SharePointSyncError(f"Chemin SharePoint invalide: {relative_path}")
    return files_dir.joinpath(*normalized_path.parts)


def _cleanup_empty_directories(root_dir: Path) -> None:
    if not root_dir.exists():
        return
    directories = [path for path in root_dir.rglob("*") if path.is_dir()]
    directories.sort(key=lambda path: (len(path.parts), path.as_posix()), reverse=True)
    for directory in directories:
        try:
            directory.rmdir()
        except OSError:
            continue


def _has_usable_cache(files_dir: Path) -> bool:
    if not files_dir.exists():
        return False
    return any(
        path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
        for path in files_dir.rglob("*")
    )


def _load_json(file_path: Path) -> dict[str, Any]:
    if not file_path.exists():
        return {}
    try:
        return json.loads(file_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}


def _write_json(file_path: Path, payload: dict[str, Any]) -> None:
    file_path.parent.mkdir(parents=True, exist_ok=True)
    file_path.write_text(json.dumps(payload, indent=2, ensure_ascii=True), encoding="utf-8")
