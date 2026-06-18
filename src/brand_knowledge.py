import json
import zipfile
from collections import Counter
from functools import lru_cache
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pypdf import PdfReader
from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
GUIDELINES_PATH = PROJECT_ROOT / "GuidelinesFR.pdf"
TOOLS_ROOT = PROJECT_ROOT / "Tools and templates PPT - FR"

OFFICIAL_TEMPLATE_FILES = [
    TOOLS_ROOT / "French" / "1. User Guide" / "OFR_Guide_Utilisateur.pptx",
    TOOLS_ROOT / "French" / "3. Best Practice" / "OFR_Best_Practice.pptx",
    TOOLS_ROOT / "French" / "4. Tools and Assets" / "OFR_Outils et elements utiles.pptx",
    TOOLS_ROOT / "French" / "2. Templates" / "French" / "OFR_template_externe.potx",
    TOOLS_ROOT / "French" / "2. Templates" / "French" / "OFR_template_confidentiel.potx",
    TOOLS_ROOT / "French" / "2. Templates" / "French" / "OFR_template_interne.potx",
    TOOLS_ROOT / "French" / "2. Templates" / "French" / "OFR_template_accessible.potx",
    TOOLS_ROOT / "French" / "2. Templates" / "French" / "OFR_template_Fond_noir.potx",
]

PDF_THEME_KEYWORDS = {
    "typographie": ["typographie", "helvetica", "police", "fonts", "font", "helvetica neue"],
    "couleurs": ["couleurs", "orange", "palette", "noir", "blanc", "color"],
    "logo": ["logo", "master logo", "small logo", "signature"],
    "photographie": ["photographie", "illustrations", "photo"],
    "digital": ["digital", "interfaces", "web"],
}


def _safe_color_string(value: Any) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    if not text or text == "None":
        return None
    return text if text.startswith("#") else f"#{text}"


def _pptx_font_size(run) -> float | None:
    font = getattr(run, "font", None)
    size = getattr(font, "size", None)
    if size is None:
        return None
    try:
        return round(float(size.pt), 2)
    except Exception:
        return None


def _extract_slide_style_exemplars(path: Path, limit: int = 400) -> list[dict[str, Any]]:
    try:
        prs = Presentation(str(path))
    except Exception:
        return []

    exemplars: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides):
        layout_name = getattr(slide.slide_layout, "name", None)
        for shape_idx, shape in enumerate(slide.shapes):
            placeholder_type = None
            if getattr(shape, "is_placeholder", False):
                try:
                    placeholder_type = str(shape.placeholder_format.type)
                except Exception:
                    placeholder_type = "placeholder"

            if not getattr(shape, "has_text_frame", False):
                continue

            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = " ".join((para.text or "").split()).strip()
                if not text:
                    continue
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    font = run.font
                    exemplars.append(
                        {
                            "source": str(path),
                            "source_kind": "official_template",
                            "file_name": path.name,
                            "slide_idx": slide_idx,
                            "shape_idx": shape_idx,
                            "para_idx": para_idx,
                            "layout_name": layout_name,
                            "placeholder_type": placeholder_type,
                            "shape_name": getattr(shape, "name", None),
                            "text_excerpt": text[:160],
                            "font_name": getattr(font, "name", None),
                            "font_size": _pptx_font_size(run),
                            "bold": getattr(font, "bold", None),
                            "italic": getattr(font, "italic", None),
                            "color": _safe_color_string(getattr(getattr(font, "color", None), "rgb", None)),
                            "alignment": str(getattr(para, "alignment", None)) if getattr(para, "alignment", None) is not None else None,
                        }
                    )
                    break
                if len(exemplars) >= limit:
                    return exemplars
    return exemplars


def _extract_theme_from_potx(path: Path) -> dict[str, Any] | None:
    if not path.exists():
        return None

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    try:
        with zipfile.ZipFile(path) as zf:
            theme = ET.fromstring(zf.read("ppt/theme/theme1.xml"))
    except Exception:
        return None

    font_scheme = theme.find(".//a:fontScheme", ns)
    major = font_scheme.find("./a:majorFont/a:latin", ns) if font_scheme is not None else None
    minor = font_scheme.find("./a:minorFont/a:latin", ns) if font_scheme is not None else None

    colors: dict[str, str] = {}
    clr_scheme = theme.find(".//a:clrScheme", ns)
    if clr_scheme is not None:
        for child in list(clr_scheme):
            first = list(child)
            if not first:
                continue
            value = first[0].attrib.get("val") or first[0].attrib.get("lastClr")
            if value:
                colors[child.tag.split("}")[-1]] = _safe_color_string(value) or value

    return {
        "path": str(path),
        "file_name": path.name,
        "major_font": major.attrib.get("typeface") if major is not None else None,
        "minor_font": minor.attrib.get("typeface") if minor is not None else None,
        "colors": colors,
    }


def _parse_xml_int(value: str | None) -> int | None:
    if not value:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _emu_to_points(value: int | None) -> float | None:
    if value is None:
        return None
    return round(value / 12700.0, 2)


def _normalize_alignment(value: str | None) -> str | None:
    mapping = {
        "l": "left",
        "ctr": "center",
        "r": "right",
        "just": "justify",
        "dist": "justify",
    }
    if not value:
        return None
    return mapping.get(value, value)


def _extract_master_text_styles_from_potx(path: Path) -> dict[str, dict[str, Any]]:
    if not path.exists():
        return {}

    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    try:
        with zipfile.ZipFile(path) as zf:
            master = ET.fromstring(zf.read("ppt/slideMasters/slideMaster1.xml"))
    except Exception:
        return {}

    tx_styles = master.find(".//p:txStyles", ns)
    if tx_styles is None:
        return {}

    style_map: dict[str, dict[str, Any]] = {}
    style_nodes = {
        "CENTER_TITLE (3)": ("titleStyle", "lvl1pPr"),
        "TITLE (1)": ("titleStyle", "lvl1pPr"),
        "SUBTITLE (4)": ("otherStyle", "lvl1pPr"),
        "BODY (2)": ("bodyStyle", "lvl1pPr"),
        "OBJECT (7)": ("bodyStyle", "lvl1pPr"),
        "DATE (16)": ("otherStyle", "lvl1pPr"),
        "FOOTER (15)": ("otherStyle", "lvl1pPr"),
        "SLIDE_NUMBER (13)": ("otherStyle", "lvl1pPr"),
    }

    for placeholder_type, (style_name, level_name) in style_nodes.items():
        style = tx_styles.find(f"p:{style_name}", ns)
        if style is None:
            continue
        level_node = style.find(f"a:{level_name}", ns)
        if level_node is None:
            continue
        default_run = level_node.find("a:defRPr", ns)
        style_map[placeholder_type] = {
            "placeholder_type": placeholder_type,
            "font_size": round((_parse_xml_int(default_run.attrib.get("sz")) or 0) / 100.0, 2) if default_run is not None and default_run.attrib.get("sz") else None,
            "bold": default_run.attrib.get("b") == "1" if default_run is not None and default_run.attrib.get("b") is not None else None,
            "italic": default_run.attrib.get("i") == "1" if default_run is not None and default_run.attrib.get("i") is not None else None,
            "alignment": _normalize_alignment(level_node.attrib.get("algn")),
            "margin_left": _emu_to_points(_parse_xml_int(level_node.attrib.get("marL"))),
            "indent": _emu_to_points(_parse_xml_int(level_node.attrib.get("indent"))),
        }

    return style_map


def _extract_guidelines_rules() -> dict[str, Any]:
    if not GUIDELINES_PATH.exists():
        return {
            "path": str(GUIDELINES_PATH),
            "available": False,
            "summary": [],
            "keywords": {},
            "priority_rules": [],
        }

    try:
        reader = PdfReader(str(GUIDELINES_PATH))
        page_texts = [(page.extract_text() or "").replace("\n", " ") for page in reader.pages]
    except Exception:
        page_texts = []

    matches: dict[str, list[str]] = {}
    for topic, keywords in PDF_THEME_KEYWORDS.items():
        topic_matches: list[str] = []
        for text in page_texts:
            lower = text.lower()
            if any(keyword in lower for keyword in keywords):
                cleaned = " ".join(text.split())
                if cleaned:
                    topic_matches.append(cleaned[:600])
            if len(topic_matches) >= 3:
                break
        matches[topic] = topic_matches

    summary = [
        "Le brandguide Orange est la source prioritaire de conformite visuelle.",
        "Le respect des regles de typographie, couleurs, logo et hierarchie visuelle doit primer sur les references metier.",
        "Le logo doit etre utilise depuis des fichiers sources officiels et positionne dans un coin avec sa zone de protection.",
        "Les presentations doivent privilegier les templates et mises en page officiels Orange FR.",
        "La forme peut changer, mais le contenu metier ne doit pas etre altere.",
    ]

    priority_rules = [
        "Toujours privilegier les templates et assets officiels Orange FR avant toute reference metier.",
        "Conserver le texte, les images et les tableaux du document source ; ne modifier que la forme.",
        "Utiliser les polices et couleurs dominantes extraites des templates officiels.",
        "Evaluer la conformite d'abord contre le brandguide, puis contre les templates officiels, puis contre le BrandStore metier.",
    ]

    return {
        "path": str(GUIDELINES_PATH),
        "available": bool(page_texts),
        "page_count": len(page_texts),
        "summary": summary,
        "keywords": matches,
        "priority_rules": priority_rules,
    }


def _build_official_templates_knowledge() -> dict[str, Any]:
    template_paths = [path for path in OFFICIAL_TEMPLATE_FILES if path.exists()]
    deck_paths = [path for path in template_paths if path.suffix.lower() in {".pptx", ".potx"}]

    exemplars: list[dict[str, Any]] = []
    layout_counter: Counter[str] = Counter()
    placeholder_counter: Counter[str] = Counter()
    font_counter: Counter[str] = Counter()
    size_counter: Counter[float] = Counter()
    color_counter: Counter[str] = Counter()
    theme_info: list[dict[str, Any]] = []
    master_text_styles: dict[str, dict[str, Any]] = {}

    for path in template_paths:
        if path.suffix.lower() == ".potx":
            theme = _extract_theme_from_potx(path)
            if theme:
                theme_info.append(theme)
            if not master_text_styles:
                master_text_styles = _extract_master_text_styles_from_potx(path)
            slide_exemplars = _extract_slide_style_exemplars(path, limit=180)
        else:
            slide_exemplars = _extract_slide_style_exemplars(path, limit=220)

        for exemplar in slide_exemplars:
            exemplars.append(exemplar)
            if exemplar.get("layout_name"):
                layout_counter[exemplar["layout_name"]] += 1
            if exemplar.get("placeholder_type"):
                placeholder_counter[exemplar["placeholder_type"]] += 1
            if exemplar.get("font_name"):
                font_counter[exemplar["font_name"]] += 1
            if exemplar.get("font_size") is not None:
                size_counter[exemplar["font_size"]] += 1
            if exemplar.get("color"):
                color_counter[exemplar["color"]] += 1

    allowed_colors: list[str] = []
    for theme in theme_info:
        for value in theme.get("colors", {}).values():
            safe = _safe_color_string(value)
            if safe and safe not in allowed_colors:
                allowed_colors.append(safe)
    for color, _ in color_counter.most_common(10):
        if color not in allowed_colors:
            allowed_colors.append(color)

    preferred_fonts = [font for font, _ in font_counter.most_common(6)]
    preferred_sizes = [size for size, _ in size_counter.most_common(12)]
    preferred_layouts = [layout for layout, _ in layout_counter.most_common(10)]
    preferred_placeholders = [item for item, _ in placeholder_counter.most_common(10)]

    return {
        "available": bool(template_paths),
        "template_paths": [str(path) for path in template_paths],
        "template_decks": [str(path) for path in deck_paths],
        "theme_info": theme_info,
        "style_exemplars": exemplars,
        "master_text_styles": master_text_styles,
        "preferred_fonts": preferred_fonts,
        "preferred_sizes": preferred_sizes,
        "allowed_colors": allowed_colors,
        "preferred_layouts": preferred_layouts,
        "preferred_placeholders": preferred_placeholders,
    }


@lru_cache(maxsize=1)
def get_brand_knowledge() -> dict[str, Any]:
    guidelines = _extract_guidelines_rules()
    templates = _build_official_templates_knowledge()
    return {
        "guidelines": guidelines,
        "templates": templates,
    }


def build_brand_reference_examples(file_type: str) -> list[dict[str, Any]]:
    if file_type != "pptx":
        return []

    knowledge = get_brand_knowledge()
    templates = knowledge["templates"]
    guidelines = knowledge["guidelines"]
    if not templates.get("available"):
        return []

    return [
        {
            "content": "\n".join(guidelines.get("summary", [])),
            "metadata": {
                "source": str(GUIDELINES_PATH),
                "filename": GUIDELINES_PATH.name,
                "extension": ".pdf",
                "priority": "brandguide",
            },
            "style_exemplars": templates.get("style_exemplars", []),
        }
    ]


def build_brand_prompt_payload(file_type: str) -> dict[str, Any]:
    knowledge = get_brand_knowledge()
    guidelines = knowledge["guidelines"]
    templates = knowledge["templates"]
    return {
        "file_type": file_type,
        "brandguide": {
            "available": guidelines.get("available"),
            "summary": guidelines.get("summary", []),
            "priority_rules": guidelines.get("priority_rules", []),
            "keywords": guidelines.get("keywords", {}),
        },
        "official_templates": {
            "available": templates.get("available"),
            "preferred_fonts": templates.get("preferred_fonts", []),
            "preferred_sizes": templates.get("preferred_sizes", []),
            "allowed_colors": templates.get("allowed_colors", []),
            "preferred_layouts": templates.get("preferred_layouts", []),
            "preferred_placeholders": templates.get("preferred_placeholders", []),
            "master_text_styles": templates.get("master_text_styles", {}),
            "template_paths": templates.get("template_paths", []),
        },
    }


def build_brand_quality_targets(file_type: str) -> dict[str, Any]:
    payload = build_brand_prompt_payload(file_type)
    templates = payload["official_templates"]
    guidelines = payload["brandguide"]
    return {
        "brandguide_available": guidelines.get("available", False),
        "priority_rules": guidelines.get("priority_rules", []),
        "preferred_fonts": templates.get("preferred_fonts", []),
        "preferred_sizes": templates.get("preferred_sizes", []),
        "allowed_colors": templates.get("allowed_colors", []),
        "preferred_layouts": templates.get("preferred_layouts", []),
        "preferred_placeholders": templates.get("preferred_placeholders", []),
        "master_text_styles": templates.get("master_text_styles", {}),
    }


def get_brand_placeholder_style(placeholder_type: str | None, file_type: str) -> dict[str, Any]:
    if file_type != "pptx":
        return {}
    targets = build_brand_quality_targets(file_type)
    master_styles = targets.get("master_text_styles") or {}
    preferred_fonts = targets.get("preferred_fonts") or []
    allowed_colors = targets.get("allowed_colors") or []
    default_style = dict(master_styles.get(placeholder_type or "") or {})

    if placeholder_type in {"CENTER_TITLE (3)", "TITLE (1)"}:
        default_style.setdefault("font_name", "Helvetica 75 Bold")
        default_style.setdefault("font_size", 20.0 if placeholder_type == "TITLE (1)" else 24.0)
        default_style.setdefault("color", "#FF7900")
        default_style.setdefault("bold", True)
    elif placeholder_type == "BODY (2)":
        default_style.setdefault("font_name", preferred_fonts[0] if preferred_fonts else "Helvetica 55 Roman")
        default_style.setdefault("font_size", 14.0)
        default_style.setdefault("color", "#000000")
    elif placeholder_type in {"DATE (16)", "FOOTER (15)", "SLIDE_NUMBER (13)", "SUBTITLE (4)"}:
        default_style.setdefault("font_name", preferred_fonts[0] if preferred_fonts else "Helvetica 55 Roman")
        default_style.setdefault("font_size", 18.0)
        default_style.setdefault("color", "#000000")

    if not default_style.get("font_name") and preferred_fonts:
        default_style["font_name"] = preferred_fonts[0]
    if not default_style.get("color") and allowed_colors:
        default_style["color"] = allowed_colors[0]

    return default_style


def get_official_template_decks(file_type: str) -> list[str]:
    if file_type != "pptx":
        return []
    return list(get_brand_knowledge()["templates"].get("template_decks", []))
