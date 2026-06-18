import json
import os
import re
import sys

from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")

from src.graph import create_graph


def _initial_state(file_path: str) -> dict:
    return {
        "file_path": file_path,
        "original_file_path": file_path,
        "file_type": None,
        "raw_elements": [],
        "anonymized_elements": [],
        "document_title": None,
        "similar_examples": [],
        "primary_reference_deck": None,
        "brand_context": None,
        "style_map": [],
        "output_path": None,
        "quality_score": None,
        "quality_report": None,
        "validation_status": "pending",
        "human_feedback": None,
        "iteration_count": 0,
        "max_iterations": int(os.getenv("MAX_ITERATIONS", "3")),
        "error": None,
        "messages": [],
    }


def _collect_pptx_style_sample(file_path: str) -> dict:
    prs = Presentation(file_path)
    sample = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                run = next((run for run in para.runs if run.text.strip()), None)
                font = getattr(run, "font", None) if run is not None else None
                color = getattr(getattr(font, "color", None), "rgb", None)
                size = getattr(font, "size", None)
                sample.append(
                    {
                        "slide_idx": slide_idx,
                        "shape_idx": shape_idx,
                        "para_idx": para_idx,
                        "text": text[:120],
                        "font_name": getattr(font, "name", None),
                        "font_size": round(float(size.pt), 2) if size is not None else None,
                        "bold": getattr(font, "bold", None),
                        "italic": getattr(font, "italic", None),
                        "color": f"#{str(color)}" if color is not None else None,
                    }
                )
                if len(sample) >= 8:
                    return {"sample": sample}

    return {"sample": sample}


def _collect_pptx_structure(file_path: str) -> dict:
    prs = Presentation(file_path)
    slides = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_items = []
        for shape_idx, shape in enumerate(slide.shapes):
            text = getattr(shape, "text", "")
            if isinstance(text, str):
                text = " ".join(text.split())[:160]
            else:
                text = ""
            placeholder = None
            if getattr(shape, "is_placeholder", False):
                try:
                    placeholder = str(shape.placeholder_format.type)
                except Exception:
                    placeholder = "placeholder"
            slide_items.append(
                {
                    "shape_idx": shape_idx,
                    "shape_type": str(shape.shape_type),
                    "placeholder": placeholder,
                    "name": getattr(shape, "name", None),
                    "text": text,
                }
            )
        slides.append({"slide_idx": slide_idx, "shape_count": len(slide_items), "items": slide_items[:10]})
    return {"slide_count": len(slides), "slides": slides[:6]}


def _compare_deck_text_similarity(output_path: str, reference_path: str) -> dict:
    def collect(path: str) -> list[str]:
        prs = Presentation(path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    lines.append(" ".join(text.split()))
        return lines

    output_lines = collect(output_path)
    reference_lines = collect(reference_path)
    output_tokens = set(re.findall(r"[a-z0-9]+", "\n".join(output_lines).lower()))
    reference_tokens = set(re.findall(r"[a-z0-9]+", "\n".join(reference_lines).lower()))
    overlap = len(output_tokens & reference_tokens)
    union = len(output_tokens | reference_tokens) or 1
    return {
        "reference_path": reference_path,
        "token_overlap_ratio": round(overlap / union, 4),
        "output_slide_count": len(Presentation(output_path).slides),
        "reference_slide_count": len(Presentation(reference_path).slides),
    }


def _extract_output_text(file_path: str) -> list[str]:
    prs = Presentation(file_path)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                lines.append(" ".join(text.split()))
    return lines


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: python scripts/run_validation.py <file_path>")
        return 1

    file_path = sys.argv[1]
    graph = create_graph()
    config = {"configurable": {"thread_id": "validation-run"}}
    last = _initial_state(file_path)

    try:
        for event in graph.stream(last, config, stream_mode="values"):
            last = event

        if last.get("error"):
            print(json.dumps({"error": last["error"]}, ensure_ascii=False, indent=2))
            return 1

        graph.update_state(config, {"validation_status": "approved", "human_feedback": None})
        for event in graph.stream(None, config, stream_mode="values"):
            last = event

        if last.get("error"):
            print(json.dumps({"error": last["error"]}, ensure_ascii=False, indent=2))
            return 1

        output_path = last.get("output_path")
        result = {
            "document_title": last.get("document_title"),
            "output_path": output_path,
            "iteration_count": last.get("iteration_count"),
            "quality_score": last.get("quality_score"),
            "style_map_count": len(last.get("style_map") or []),
            "raw_element_count": len(last.get("raw_elements") or []),
            "similar_examples_count": len(last.get("similar_examples") or []),
            "primary_reference_deck": last.get("primary_reference_deck"),
            "brand_context": last.get("brand_context"),
            "pptx_style_sample": _collect_pptx_style_sample(output_path) if output_path and output_path.lower().endswith(".pptx") else None,
            "pptx_structure": _collect_pptx_structure(output_path) if output_path and output_path.lower().endswith(".pptx") else None,
        }
        if output_path and output_path.lower().endswith(".pptx") and last.get("primary_reference_deck"):
            result["reference_comparison"] = _compare_deck_text_similarity(
                output_path,
                last["primary_reference_deck"]["path"],
            )
        if output_path and output_path.lower().endswith(".pptx"):
            result["output_text_sample"] = _extract_output_text(output_path)[:20]
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    except Exception as exc:
        print(json.dumps({"error": str(exc)}, ensure_ascii=False, indent=2))
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
